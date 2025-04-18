import argparse
import os
import re
from datetime import datetime
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
import requests
import keyring
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

TITLE_IDX = 0
DESCRIPTION_IDX = 1
IMAGE_IDX = 10
PB_LINK_IDX = 11
ID_IDX = 12
REQUIREMENTS_IDX = 13

INCLUDED_STATUS_IDS = [
    #"402b0df5-1554-44bf-bcc4-8a11d3ca0a65",  # Candidate
    "b944953f-2a93-4e85-b036-9bef92432588", #Planned
    "d3dad2ee-2692-4e06-a29c-4ec314f807a0", #In Progress
    "d6c9c0da-411a-41a9-8343-a8a5f75036e3", #EG.EG
    "93d3820a-af3f-43e3-9122-32e302d7efd1", #Beta
    "65a03095-72c0-42f2-952a-ab1c30abdae4", #Limited Release 
    "88c3f59c-2cb0-4438-8840-d00f85fce6d1" # Released
]

EXCLUDED_STATUS_IDS = [
    "402b0df5-1554-44bf-bcc4-8a11d3ca0a65",  # Candidate
    "310a6c38-3719-4bff-a26c-b5d520ea1298",  # New Idea
]
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

api_token = keyring.get_password("productboard-api", "default")
if not api_token:
    raise RuntimeError("API token not found in keyring.")

HEADERS = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {api_token}"
}

TIMEFRAME_START = datetime.fromisoformat("2025-04-01")
TIMEFRAME_END = datetime.fromisoformat("2025-06-30")

#################### IMAGE-HANDLING FUNCTIONS ####################

def extract_image_urls(description_html):
    
    # Finds <img src='...'> in the HTML description and returns them, also removing those tags from the HTML so it doesn't appear as broken text.
    
    img_tags = re.findall(r'<img [^>]*src="([^"]+)"[^>]*>', description_html)
    cleaned_html = re.sub(r'<img [^>]*src="[^"]+"[^>]*>', '', description_html)
    return img_tags, cleaned_html

def fetch_image(url):
    # Fetches and verifies an image if it's hosted on pb-files.s3.amazonaws.com, if this function is not added, Emojis will be added to all features.
    try:
        if "pb-files.s3.amazonaws.com" not in url:
            print(f"Skipping non-Productboard image URL: {url}")
            return None

        resp = requests.get(url)
        if resp.status_code != 200:
            print(f"Failed to fetch image from {url}. Status code: {resp.status_code}")
            return None

        content_type = resp.headers.get("Content-Type")
        if not content_type or not content_type.startswith("image/"):
            print(f"URL is not an image: {url}")
            return None

        # Verify image integrity
        img = Image.open(BytesIO(resp.content))
        img.verify()
        return Image.open(BytesIO(resp.content))  # re-load
    except Exception as e:
        print(f"Error fetching or processing image from {url}: {e}")
        return None

def insert_image_with_aspect_ratio(slide, placeholder, img):
    
    # Inserts a PIL image into the 'placeholder' shape while preserving aspect ratio, centered within the placeholder's bounding box.

    ph_width = placeholder.width
    ph_height = placeholder.height

    img_width, img_height = img.size
    ph_aspect = ph_width / ph_height
    img_aspect = img_width / img_height

    if img_aspect > ph_aspect:
        # Width-limited
        new_width = ph_width
        new_height = int(ph_width / img_aspect)
    else:
        # Height-limited
        new_height = ph_height
        new_width = int(ph_height * img_aspect)

    left = placeholder.left + int((ph_width - new_width) / 2)
    top = placeholder.top + int((ph_height - new_height) / 2)

    img_stream = BytesIO()
    img.save(img_stream, format="PNG")
    img_stream.seek(0)

    slide.shapes.add_picture(img_stream, left, top, width=new_width, height=new_height)
    print("✅ Inserted image with aspect ratio and centered.")

#################### TEXT HELPER FUNCTIONS ####################

def safe_get(value, default="Missing requirements"):
    if value is None:
        return default
    return str(value).strip() if str(value).strip() else default

def fill_empty_text_if_needed(shape):
    """If a placeholder is truly empty, fill with a space to avoid PPTX corruption."""
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text.strip()
    if not text:
        shape.text_frame.text = " "

def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
    """Used in HTML formatting or anywhere we build runs."""
    if not text.strip():
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.font.bold = bold
    run.font.underline = underline
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

def clean_html_and_format_text(description_html, text_frame):
    """Clears the text frame and formats HTML into PPT paragraphs,
    fully supporting arbitrarily nested <ul>/<ol> lists."""
    
    # 1) Clear out any existing paragraphs
    while text_frame.paragraphs:
        p = text_frame.paragraphs[-1]
        for r in p.runs:
            r._r.getparent().remove(r._r)
        p._element.getparent().remove(p._element)

    soup = BeautifulSoup(description_html or "", "html.parser")
    if not soup.contents:
        p = text_frame.add_paragraph()
        add_run(p, "The description is empty.")
        return

    def render_list(list_tag, level=0):
        ordered = (list_tag.name == "ol")
        for idx, li in enumerate(list_tag.find_all("li", recursive=False), start=1):
            # build the line text (skip nested lists)
            line = ""
            nested = []
            for c in li.contents:
                if isinstance(c, NavigableString):
                    line += c.strip()
                elif isinstance(c, Tag) and c.name not in ("ul", "ol"):
                    line += c.get_text(strip=True)
                elif isinstance(c, Tag) and c.name in ("ul","ol"):
                    nested.append(c)

            # emit paragraph
            p = text_frame.add_paragraph()
            p.level = level
            prefix = f"{idx}. " if ordered else "• "
            add_run(p, prefix + line)

            # recurse into each nested list
            for nl in nested:
                render_list(nl, level + 1)

    # walk top‑level blocks
    for block in soup.find_all(recursive=False):
        if isinstance(block, NavigableString):
            text = block.strip()
            if text:
                p = text_frame.add_paragraph()
                add_run(p, text)

        elif isinstance(block, Tag):
            if block.name in ("h1","h2","h3","h4"):
                p = text_frame.add_paragraph()
                add_run(p, block.get_text(strip=True), bold=True)
                p.space_after = Pt(8)

            elif block.name == "p":
                content = block.get_text(strip=True)
                if content:
                    p = text_frame.add_paragraph()
                    for child in block.children:
                        if isinstance(child, NavigableString):
                            add_run(p, child.strip())
                        elif child.name in ("strong","b"):
                            add_run(p, child.get_text(strip=True), bold=True)
                        elif child.name == "a":
                            add_run(p,
                                    child.get_text(strip=True),
                                    underline=True,
                                    hyperlink=child.get("href",""))
                    p.space_after = Pt(6)

            elif block.name in ("ul","ol"):
                render_list(block, level=0)

    # strip empty paragraphs
    for p in list(text_frame.paragraphs):
        if not p.text.strip():
            e = p._element
            e.getparent().remove(e)

    print("✅ Finished formatting text (nested lists supported)")

#################### DATA FETCH & FILTERING ####################

def get_all_paginated_features(url):
    all_features = []
    page = 1
    while url:
        print(f"📅 Fetching features page {page}: {url}")
        response = requests.get(url, headers=HEADERS)
        if response.status_code != 200:
            print(f"❌ Failed to fetch features: {response.text}")
            break
        data = response.json()
        features = data.get("data", [])
        print(f"🔹 Retrieved {len(features)} features on page {page}")
        all_features.extend(features)
        url = data.get("links", {}).get("next")
        page += 1
    print(f"📆 Total features retrieved: {len(all_features)}")
    return all_features

def get_feature_ids_by_status_id(status_id):
    url = f"https://api.productboard.com/features?status.id={status_id}"
    return set(f["id"] for f in get_all_paginated_features(url))

def get_feature_details(fid):
    url = f"https://api.productboard.com/features/{fid}"
    response = requests.get(url, headers=HEADERS)
    print(f"📄 Feature {fid} details status: {response.status_code}")
    return response.json()

def get_requirements_link(fid):
    custom_field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={custom_field_id}&hierarchyEntity.id={fid}"
    response = requests.get(url, headers=HEADERS)
    return response.json().get("data", {}).get("value", None)

def get_jira_details(fid):
    url = f"https://api.productboard.com/jira-integrations/{JIRA_API_ID}/connections/{fid}"
    response = requests.get(url, headers=HEADERS)
    if response.status_code != 200:
        print(f"⚠️ Jira link not found for {fid}: {response.status_code}")
        return None, None
    jira_data = response.json()
    issue_key = jira_data.get("data", {}).get("connection", {}).get("issueKey")
    if issue_key:
        return issue_key, f"https://jira.egnyte-it.com/browse/{issue_key}"
    return None, None

def get_placeholder_by_idx(slide, idx):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            # Fill any truly empty text with a single space
            if shape.has_text_frame and not shape.text_frame.text.strip():
                shape.text_frame.text = " "
            return shape
    return None

#################### SAFE TEXT FUNCTION ####################

def safe_clear_and_add_text(text_frame, text, hyperlink=None, size=Pt(11), font_name="Avenir"):
    """
    Clears all paragraphs in text_frame, then adds exactly one paragraph/run with 'text'.
    If text is empty, fallback to " " to avoid PPTX corruption (empty <a:p/>).
    """
    while text_frame.paragraphs:
        paragraph = text_frame.paragraphs[-1]
        for r in paragraph.runs:
            r._r.getparent().remove(r._r)  # Use r._r for the underlying <a:r>
        paragraph._element.getparent().remove(paragraph._element)

    paragraph = text_frame.add_paragraph()
    text = text.strip() if text else ""
    run = paragraph.add_run()
    run.text = text if text else " "
    run.font.size = size
    run.font.name = font_name

    # Only attach hyperlink if we have a nonempty text
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

#################### CREATE PPTX ####################

def create_pptx(features_data):
    prs = Presentation("templates/corporate_template.pptx")

    # Group features by 'initiative'
    grouped = {}
    for f in features_data:
        grouped.setdefault(f["initiative"], []).append(f)

    # For each initiative, create a cover slide and feature slides
    for initiative, items in grouped.items():
        cover_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(cover_layout)
        slide.shapes.title.text = initiative

        for feature in items:
            # 1) Extract images from the description
            all_images, cleaned_html = extract_image_urls(feature.get("description", "") or "")

            # 2) Create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = feature['title']

            # 3) Populate description
            desc_placeholder = get_placeholder_by_idx(slide, DESCRIPTION_IDX)
            if desc_placeholder and desc_placeholder.has_text_frame:
                clean_html_and_format_text(cleaned_html, desc_placeholder.text_frame)

            # 4) Insert images if we have them
            if all_images:
                img_placeholder = get_placeholder_by_idx(slide, IMAGE_IDX)
                if img_placeholder:
                    # Insert up to 4 images
                    for i, url in enumerate(all_images[:4], 1):
                        print(f"🔍 Fetching image {i} for feature: {feature['title']} => {url}")
                        pil_img = fetch_image(url)
                        if pil_img:
                            insert_image_with_aspect_ratio(slide, img_placeholder, pil_img)
                        else:
                            print(f"⚠️ Skipped invalid image at {url}")

            # 5) Productboard link
            pb_placeholder = get_placeholder_by_idx(slide, PB_LINK_IDX)
            if pb_placeholder and pb_placeholder.has_text_frame:
                link = feature.get("html_link")
                safe_clear_and_add_text(pb_placeholder.text_frame, "View in Productboard" if link else "", hyperlink=link)

            # 6) JIRA link
            jira_placeholder = get_placeholder_by_idx(slide, ID_IDX)
            if jira_placeholder and jira_placeholder.has_text_frame:
                jira_key = feature.get("jira_key")
                jira_url = feature.get("jira_url")
                if jira_key and jira_url:
                    safe_clear_and_add_text(jira_placeholder.text_frame, jira_key, hyperlink=jira_url)
                else:
                    safe_clear_and_add_text(jira_placeholder.text_frame, "No Jira Link")

            # 7) Requirements
            req_placeholder = get_placeholder_by_idx(slide, REQUIREMENTS_IDX)
            if req_placeholder and req_placeholder.has_text_frame:
                requirements_link = (feature.get("requirements_link") or "").strip()
                if requirements_link:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Link to requirements", hyperlink=requirements_link)
                else:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Missing requirements")

            # 8) Final housekeeping
            for shape in slide.placeholders:
                fill_empty_text_if_needed(shape)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    prs.save(f"output_presentation_{timestamp}.pptx")
    print(f"✅ Presentation saved as output_presentation_{timestamp}.pptx")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--owner_email", default=None, help="Filter by owner email")
    args = parser.parse_args()

    # 1) Fetch only the statuses you care about
    print("🔍 Fetching only specified statuses…")
    feature_ids = set()
    for status_id in INCLUDED_STATUS_IDS:
        print(f"🔎 Fetching features with status {status_id}")
        feature_ids |= get_feature_ids_by_status_id(status_id)

    # 2) (Optional) preserve the API’s original ordering
    print("📋 Retrieving all features to preserve order…")
    all_features = get_all_paginated_features("https://api.productboard.com/features")
    remaining_ids = [f["id"] for f in all_features if f["id"] in feature_ids]
    print(f"✅ Will process {len(remaining_ids)} features\n")

    # 3) Build initiative lookups
    print("📂 Fetching initiatives…")
    initiatives = get_all_paginated_features("https://api.productboard.com/initiatives")
    initiative_map = {i["id"]: i["name"] for i in initiatives}

    print("🔗 Building feature→initiative map…")
    feature_to_initiative = {}
    for iid in initiative_map:
        links = get_all_paginated_features(
            f"https://api.productboard.com/initiatives/{iid}/links/features"
        )
        for link in links:
            fid = link.get("id")
            if fid:
                feature_to_initiative[fid] = iid

    # 4) Fetch full feature details in parallel, filter by timeframe + owner
    features = []
    print("🧵 Fetching feature details in parallel…")
    with ThreadPoolExecutor(max_workers=20) as executor:
        futures = {
            executor.submit(get_feature_details, fid): fid
            for fid in remaining_ids
        }
        for future in as_completed(futures):
            fid = futures[future]
            try:
                details = future.result().get("data", {})
                tf = details.get("timeframe", {})
                start, end = tf.get("startDate"), tf.get("endDate")
                owner = details.get("owner") or {}
                email = owner.get("email", "")

                # Skip if no valid dates
                if not (start and end):
                    continue

                s_dt = datetime.fromisoformat(start)
                e_dt = datetime.fromisoformat(end)
                if not (s_dt <= TIMEFRAME_END and e_dt >= TIMEFRAME_START):
                    continue

                # Skip if owner_email filter is set
                if args.owner_email and args.owner_email != email:
                    continue

                # Gather links
                req_link = get_requirements_link(fid)
                jira_key, jira_url = get_jira_details(fid)
                initiative = initiative_map.get(
                    feature_to_initiative.get(fid), "Uncategorized"
                )

                features.append({
                    "id": fid,
                    "title": details.get("name", ""),
                    "description": details.get("description", ""),
                    "requirements_link": req_link,
                    "html_link": details.get("links", {}).get("html", ""),
                    "initiative": initiative,
                    "jira_key": jira_key,
                    "jira_url": jira_url
                })

            except Exception as e:
                print(f"❌ Error fetching feature {fid}: {e}")

    print(f"\n📊 Final features to generate slides for: {len(features)}\n")

    # 5) Create the PPTX
    create_pptx(features)


if __name__ == "__main__":
    main()

