import argparse
import re
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import requests
import keyring
from bs4 import BeautifulSoup, NavigableString, Tag

TITLE_IDX = 0
DESCRIPTION_IDX = 1
IMAGE_IDX = 10
PB_LINK_IDX = 11
ID_IDX = 12
REQUIREMENTS_IDX = 13
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

INITIATIVE_FEATURE_LIST_IDX = 1
INITIATIVE_DESC_IDX = 10
INITIATIVE_EMPTY_IDS = [11,12,13]

TIMEFRAME_START = datetime.fromisoformat("2025-04-01")
TIMEFRAME_END = datetime.fromisoformat("2025-06-30")

initiative_ids = [
    "5f9c6433-028a-4779-b50a-fa8be3d35ff7", #Introduce Gen 4 Plans
    "bc714ff2-5193-403d-9be8-cd7dd59f165f", #Top Customer Requests
    "0305631a-607f-48b7-9f0a-16e88e99d721", #Elevate Desktop App Capabilities
    "16366226-bcbc-4301-ade4-2fb7bbcd0cb3", #Productivity Agents
    "03ea3159-fe6f-4b87-97ac-0dea07b681b0", #Specialized Agents
    "b47dd248-f4d4-4374-9b67-04e5faf1fbd7", #Augment AEC Add-on Value
    "9370064c-5fac-432a-ad53-21a72fad007a", #LS Mid Market
    "347e0a98-9efe-44d7-91e2-a712283a2d8f", #Co-editing
    "48e537ad-513c-4afe-90cf-55f35a2ea6ec", #Product-Led Grwoth
    "21e6df72-36d0-4749-8612-c59f8dfc1ff9" #doc portal enhancements
]

api_token = keyring.get_password("productboard-api", "default")
HEADERS = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {api_token}"
}
def get_feature_details(feature):
#   Fetch detailed information for a feature.   If the feature dictionary includes a 'self_link', that URL is used to request the details. Otherwise, it falls back to constructing the URL using the feature's id.
    
    if "self_link" in feature and feature["self_link"]:
        url = feature["self_link"]
    else:
        url = f"https://api.productboard.com/features/{feature.get('id')}"
    response = requests.get(url, headers=HEADERS)
    print(f"📄 Feature {feature.get('id')} details status: {response.status_code}")
    return response.json()


def get_linked_features(initiative_id):
    url = f"https://api.productboard.com/initiatives/{initiative_id}/links/features"
    return requests.get(url, headers=HEADERS).json().get("data", [])

def get_feature_details(fid):
    url = f"https://api.productboard.com/features/{fid}"
    return requests.get(url, headers=HEADERS).json().get("data", {})

def get_requirements_link(fid):
    field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={field_id}&hierarchyEntity.id={fid}"
    return requests.get(url, headers=HEADERS).json().get("data", {}).get("value")

def get_jira_details(fid):
    url = f"https://api.productboard.com/jira-integrations/{JIRA_API_ID}/connections/{fid}"
    r = requests.get(url, headers=HEADERS)
    if r.status_code != 200:
        return None, None
    d = r.json().get("data", {}).get("connection", {})
    return d.get("issueKey"), f"https://jira.egnyte-it.com/browse/{d.get('issueKey')}" if d.get("issueKey") else None

def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
    if not text.strip():
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.font.bold = bold
    run.font.underline = underline
    if hyperlink and text:
        run.hyperlink.address = hyperlink

def clean_html_and_format_text(html, text_frame):
    soup = BeautifulSoup(html or "", "html.parser")
    text_frame.clear()

    for element in soup.contents:
        p = text_frame.add_paragraph()
        if isinstance(element, NavigableString):
            add_run(p, element.strip())
        elif element.name == "p":
            for child in element.children:
                if isinstance(child, NavigableString):
                    add_run(p, child.strip())
                elif child.name in ["strong", "b"]:
                    add_run(p, child.get_text(strip=True), bold=True)
                elif child.name == "a":
                    href = child.get("href", "")
                    link_text = child.get_text(strip=True)
                    add_run(p, link_text, underline=True, hyperlink=href)
        elif element.name in ["ul", "ol"]:
            is_ordered = element.name == "ol"
            for idx, li in enumerate(element.find_all("li", recursive=False), 1):
                p = text_frame.add_paragraph()
                prefix = f"{idx}. " if is_ordered else "• "
                p.text = prefix + li.get_text(strip=True)

def get_placeholder_by_idx(slide, idx):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None

def safe_clear_and_add_text(text_frame, text, hyperlink=None, size=Pt(11), font_name="Avenir"):
    while text_frame.paragraphs:
        paragraph = text_frame.paragraphs[-1]
        for run in paragraph.runs:
            run._element.getparent().remove(run._element)
        paragraph._element.getparent().remove(paragraph._element)
    paragraph = text_frame.add_paragraph()
    run = paragraph.add_run()
    run.text = text or " "
    run.font.size = size
    run.font.name = font_name
    if hyperlink and text:
        run.hyperlink.address = hyperlink

INITIATIVE_LAYOUT_IDX = 2

def create_initiative_slide(prs, initiative_detail, features, feature_slide_map): """ Create an initiative slide. - The placeholder at INITIATIVE_FEATURE_LIST_IDX will contain a sorted list of features, each with an internal hyperlink to the detailed feature slide. - Placeholder INITIATIVE_DESC_IDX contains the initiative description. - Placeholders with idx 11, 12, and 13 are cleared. """ slide_layout = prs.slide_layouts[INITIATIVE_LAYOUT_IDX] slide = prs.slides.add_slide(slide_layout)


# Fill the feature list placeholder (IDX 1)
feature_list_placeholder = slide.placeholders[initiative_ids]
# Clear any existing paragraphs in the text frame
feature_list_placeholder.text_frame.clear()
# Sort the features alphabetically by title
sorted_features = sorted(features, key=lambda x: x.get("title", ""))
for feat in sorted_features:
    p = feature_list_placeholder.text_frame.add_paragraph()
    # Add the feature title text
    run = p.add_run()
    run.text = feat.get("title", "Untitled")
    # If we have a slide id mapping, add an internal hyperlink
    target_slide_id = feature_slide_map.get(feat.get("id"))
    if target_slide_id:
        add_internal_hyperlink(run, target_slide_id)
print("✅ Initiative feature list added.")

# Fill the initiative description in placeholder IDX 10
init_desc_ph = slide.placeholders[INITIATIVE_DESC_IDX]
init_desc_ph.text = initiative_detail.get("description", "")

# Clear placeholders with idx 11, 12, 13
for idx in INITIATIVE_EMPTY_IDS:
    ph = slide.placeholders[idx]
    ph.text = ""

print(f"✅ Initiative slide created for: {initiative_detail.get('name', 'Unnamed initiative')}")
return slide

#################### MAIN FUNCTION ####################

def main(): parser = argparse.ArgumentParser() parser.add_argument("--owner_email", default=None, help="Filter by owner email") args = parser.parse_args()

python
Copy
# Load the corporate template (ensure this file exists in templates/)
prs = Presentation("templates/corporate_template.pptx")
# Mapping from feature id to its slide id (for internal linking)
feature_slide_map = {}

# Process each hardcoded initiative
for initiative_id in INITIATIVE_IDS:
    print(f"\n🔍 Processing initiative: {initiative_id}")
    initiative_detail = get_initiative_details(initiative_id)
    if not initiative_detail:
        print(f"❌ Skipping initiative {initiative_id} due to missing details.")
        continue

    # Get feature IDs linked to this initiative
    feature_ids = get_features_for_initiative(initiative_id)
    initiative_features = []
    # Fetch feature details in parallel for this initiative
    with ThreadPoolExecutor(max_workers=10) as executor:
        future_to_fid = {executor.submit(get_feature_details, fid): fid for fid in feature_ids}
        for future in as_completed(future_to_fid):
            fid = future_to_fid[future]
            try:
                details = future.result()
                data = details.get("data")
                if not isinstance(data, dict):
                    continue

                # Filter features by timeframe (Q2)
                timeframe = data.get("timeframe", {})
                start, end = timeframe.get("startDate"), timeframe.get("endDate")
                if not (start and end):
                    continue
                try:
                    s_dt = datetime.fromisoformat(start)
                    e_dt = datetime.fromisoformat(end)
                except Exception:
                    print(f"⚠️ Invalid date for feature {fid}: {start} - {end}")
                    continue
                if not (s_dt <= TIMEFRAME_END and e_dt >= TIMEFRAME_START):
                    continue

                owner = data.get("owner") or {}
                email = owner.get("email", "")
                if args.owner_email and args.owner_email != email:
                    continue

                req = get_requirements_link(fid)
                jira_key, jira_url = get_jira_details(fid)
                feature_data = {
                    "id": fid,
                    "title": data.get("name", ""),
                    "description": data.get("description", ""),
                    "requirements_link": req,
                    "html_link": data.get("links", {}).get("html", ""),
                    "jira_key": jira_key,
                    "jira_url": jira_url
                }
                initiative_features.append(feature_data)
            except Exception as e:
                print(f"❌ Error processing feature {fid}: {e}")

    print(f"📊 Initiative {initiative_id} – {len(initiative_features)} features after filtering for Q2.")

    # Create detailed feature slides for all features in this initiative
    for feature in initiative_features:
        slide_id = create_feature_slide(prs, feature)
        # Map the feature id to its slide id for hyperlinking
        feature_slide_map[feature["id"]] = slide_id

    # Create the initiative slide that summarizes the features
    create_initiative_slide(prs, initiative_detail, initiative_features, feature_slide_map)

# Save the final presentation with a timestamp in the filename
timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
output_filename = f"output_presentation_{timestamp}.pptx"
prs.save(output_filename)
print(f"\n✅ Presentation saved as {output_filename}")
if name == "main": main()