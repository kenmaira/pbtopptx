from pptx.util import Pt
from pptx import Presentation
from .text_utils import fill_empty_text_if_needed, clean_html_and_format_text
from .image_utils import extract_image_urls, fetch_image, insert_image_with_aspect_ratio
from .data_fetch import get_placeholder_by_idx

# PPTX creation and safe text utilities for modular_pbtopptx

TITLE_IDX = 0
DESCRIPTION_IDX = 1
IMAGE_IDX = 10
PB_LINK_IDX = 11
ID_IDX = 12
REQUIREMENTS_IDX = 13

def safe_clear_and_add_text(text_frame, text, hyperlink=None, size=Pt(11), font_name="Avenir"):
    while text_frame.paragraphs:
        paragraph = text_frame.paragraphs[-1]
        for r in paragraph.runs:
            r._r.getparent().remove(r._r)
        paragraph._element.getparent().remove(paragraph._element)
    paragraph = text_frame.add_paragraph()
    text = text.strip() if text else ""
    run = paragraph.add_run()
    run.text = text if text else " "
    run.font.size = size
    run.font.name = font_name
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

def create_pptx(features_data, template_path="templates/corporate_template.pptx"):
    prs = Presentation(template_path)
    grouped = {}
    for f in features_data:
        grouped.setdefault(f["initiative"], []).append(f)
    for initiative, items in grouped.items():
        cover_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(cover_layout)
        slide.shapes.title.text = initiative
        for feature in items:
            all_images, cleaned_html = extract_image_urls(feature.get("description", "") or "")
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = feature['title']
            desc_placeholder = get_placeholder_by_idx(slide, DESCRIPTION_IDX)
            if desc_placeholder and desc_placeholder.has_text_frame:
                clean_html_and_format_text(cleaned_html, desc_placeholder.text_frame)
            if all_images:
                img_placeholder = get_placeholder_by_idx(slide, IMAGE_IDX)
                if img_placeholder:
                    for i, url in enumerate(all_images[:4], 1):
                        print(f"🔍 Fetching image {i} for feature: {feature['title']} => {url}")
                        pil_img = fetch_image(url)
                        if pil_img:
                            insert_image_with_aspect_ratio(slide, img_placeholder, pil_img)
                        else:
                            print(f"⚠️ Skipped invalid image at {url}")
            pb_placeholder = get_placeholder_by_idx(slide, PB_LINK_IDX)
            if pb_placeholder and pb_placeholder.has_text_frame:
                link = feature.get("html_link")
                safe_clear_and_add_text(pb_placeholder.text_frame, "View in Productboard" if link else "", hyperlink=link)
            jira_placeholder = get_placeholder_by_idx(slide, ID_IDX)
            if jira_placeholder and jira_placeholder.has_text_frame:
                jira_key = feature.get("jira_key")
                jira_url = feature.get("jira_url")
                if jira_key and jira_url:
                    safe_clear_and_add_text(jira_placeholder.text_frame, jira_key, hyperlink=jira_url)
                else:
                    safe_clear_and_add_text(jira_placeholder.text_frame, "No Jira Link")
            req_placeholder = get_placeholder_by_idx(slide, REQUIREMENTS_IDX)
            if req_placeholder and req_placeholder.has_text_frame:
                requirements_link = (feature.get("requirements_link") or "").strip()
                if requirements_link:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Link to requirements", hyperlink=requirements_link)
                else:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Missing requirements")
            for shape in slide.placeholders:
                fill_empty_text_if_needed(shape)
    from datetime import datetime
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    prs.save(f"output_presentation_{timestamp}.pptx")
    print(f"✅ Presentation saved as output_presentation_{timestamp}.pptx")
