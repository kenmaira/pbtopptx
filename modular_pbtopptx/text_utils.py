import re
from pptx.util import Pt
from bs4 import BeautifulSoup
from bs4.element import NavigableString, Tag

def safe_get(value, default="Missing requirements"):
    """Safely get a value from a potentially missing attribute."""
    if value is None:
        return default
    return str(value).strip() if str(value).strip() else default

def fill_empty_text_if_needed(shape):
    """Fill empty text in a shape with a placeholder if needed."""
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text.strip()
    if not text:
        shape.text_frame.text = " "

def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
    """Add a run to a paragraph with specified formatting."""
    if not text.strip():
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.font.bold = bold
    run.font.underline = underline
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

def clean_html_and_format_text(description_html, text_frame):
    """Clean HTML tags from the description and format the text."""
    while text_frame.paragraphs:
        p = text_frame.paragraphs[-1]
        for r in p.runs:
            r._r.getparent().remove(r._r)
        p._element.getparent().remove(p._element)
    soup = BeautifulSoup(description_html or "", "html.parser")
    if not soup.contents:
        p = text_frame.add_paragraph()
        add_run(p, "The description is empty.")
        return
    def render_list(list_tag, level=0):
        ordered = (list_tag.name == "ol")
        for idx, li in enumerate(list_tag.find_all("li", recursive=False), start=1):
            line = ""
            nested = []
            for c in li.contents:
                if isinstance(c, NavigableString):
                    line += c.strip()
                elif isinstance(c, Tag) and c.name not in ("ul", "ol"):
                    line += c.get_text(strip=True)
                elif isinstance(c, Tag) and c.name in ("ul","ol"):
                    nested.append(c)
            p = text_frame.add_paragraph()
            p.level = level
            prefix = f"{idx}. " if ordered else "• "
            add_run(p, prefix + line)
            for nl in nested:
                render_list(nl, level + 1)
    for block in soup.find_all(recursive=False):
        if isinstance(block, NavigableString):
            text = block.strip()
            if text:
                p = text_frame.add_paragraph()
                add_run(p, text)
        elif isinstance(block, Tag):
            if block.name in ("h1","h2","h3","h4"):
                p = text_frame.add_paragraph()
                add_run(p, block.get_text(strip=True), bold=True)
                p.space_after = Pt(8)
            elif block.name == "p":
                content = block.get_text(strip=True)
                if content:
                    p = text_frame.add_paragraph()
                    for child in block.children:
                        if isinstance(child, NavigableString):
                            add_run(p, child.strip())
                        elif isinstance(child, Tag) and child.name in ("strong","b"):
                            add_run(p, child.get_text(strip=True), bold=True)
                        elif isinstance(child, Tag) and child.name == "a":
                            add_run(p,
                                    child.get_text(strip=True),
                                    underline=True,
                                    hyperlink=child.get("href",""))
                    p.space_after = Pt(6)
            elif block.name in ("ul","ol"):
                render_list(block, level=0)
    for p in list(text_frame.paragraphs):
        if not p.text.strip():
            e = p._element
            e.getparent().remove(e)
    print("✅ Finished formatting text (nested lists supported)")
