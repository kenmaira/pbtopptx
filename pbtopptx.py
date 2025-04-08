import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image
from datetime import datetime
from bs4 import BeautifulSoup
import argparse
import re
import os

# Headers for API calls
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

headers = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {os.getenv('PRODUCTBOARD_API_TOKEN')}"
}
print(f"Using token: {headers['authorization']}")

# Step 1: Retrieve all features in a release
def get_feature_uuids(release_id):
    url = f"https://api.productboard.com/feature-release-assignments?release.id={release_id}"
    response = requests.get(url, headers=headers)
    print(f"API response for feature UUIDs: {response.status_code}")  # Debug: Check API response status
    features = response.json().get("data", [])
    feature_ids = [feature['feature']['id'] for feature in features]
    print(f"Feature UUIDs retrieved: {feature_ids}")  # Debug: Check the feature IDs
    return feature_ids

# Step 2: Get feature details
def get_feature_details(feature_id):
    url = f"https://api.productboard.com/features/{feature_id}"
    response = requests.get(url, headers=headers)
    print(f"API response for feature details: {response.status_code}")  # Debug: Check API response status
    return response.json()

# Step 3: Get intiative details
def get_initiatives():
    url = "https://api.productboard.com/initiatives"
    response = requests.get(url, headers=headers)
    print(f"API response for initiatives: {response.status_code}")

    if response.status_code != 200:
        print(f"Failed to fetch initiatives: {response.text}")
        return {}

    initiatives_data = response.json().get("data", [])

    # Create an initiative mapping {initiative_id: initiative_name}
    initiatives = {initiative["id"]: initiative["name"] for initiative in initiatives_data}

    print(f"Retrieved {len(initiatives)} initiatives: {initiatives}")
    return initiatives

# Step 4: Get feature release assignments
def get_feature_release_assignments(release_id):
    url = f"https://api.productboard.com/feature-release-assignments?release.id={release_id}"
    response = requests.get(url, headers=headers)
    print(f"API response for feature-release assignments: {response.status_code}")

    if response.status_code != 200:
        print(f"Failed to fetch feature-release assignments: {response.text}")
        return []

    features_data = response.json().get("data", [])
    
    if not features_data:
        print(f"⚠️ Warning: No features found for release {release_id}!")

    feature_ids = [feature['feature']['id'] for feature in features_data]
    print(f"✅ Retrieved {len(feature_ids)} features assigned to release {release_id}: {feature_ids}")
    return feature_ids

# Step 5: Cross-reference features and initiatives
def get_initiative_feature_links(initiatives):
    feature_to_initiative = {}

    for initiative_id, initiative_name in initiatives.items():
        url = f"https://api.productboard.com/initiatives/{initiative_id}/links/features"
        response = requests.get(url, headers=headers)
        print(f"API response for initiative {initiative_name} ({initiative_id}): {response.status_code}")

        if response.status_code != 200:
            print(f"⚠️ Warning: Failed to fetch features for initiative {initiative_name}: {response.text}")
            continue

        links_data = response.json().get("data", [])

        # Debug: Log raw API response
        print(f"🔎 Raw Response for Initiative {initiative_name}: {links_data}")

        if not links_data:
            print(f"⚠️ No features found for initiative {initiative_name} ({initiative_id})!")

        # Map each feature to its initiative
        for link in links_data:
            feature_id = link.get("id")
            if feature_id:
                feature_to_initiative[feature_id] = initiative_id
                print(f"✅ Feature {feature_id} linked to Initiative: {initiative_name}")

    print(f"✅ Retrieved {len(feature_to_initiative)} feature-to-initiative links")
    return feature_to_initiative

# Step 6: Combine data & create a structured output
def group_features_by_initiative(release_id):
    initiatives = get_initiatives()
    release_features = get_feature_release_assignments(release_id)

    if not release_features:
        print(f"❌ ERROR: No features assigned to release {release_id}. Check Productboard!")
        return {}

    feature_to_initiative = get_initiative_feature_links(initiatives)

    grouped_features = {}

    for feature_id in release_features:
        initiative_id = feature_to_initiative.get(feature_id, None)
        initiative_name = initiatives.get(initiative_id, "Uncategorized")

        print(f"🔎 Feature {feature_id} → Initiative ID: {initiative_id} → Initiative Name: {initiative_name}")

        if initiative_name not in grouped_features:
            grouped_features[initiative_name] = []
        
        grouped_features[initiative_name].append(feature_id)

    print(f"✅ Final Grouped Features → {grouped_features}")
    return grouped_features

# Step 7: Get custom field (requirements link)
def get_requirements_link(feature_id):
    custom_field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={custom_field_id}&hierarchyEntity.id={feature_id}"
    response = requests.get(url, headers=headers)
    print(f"API response for requirements link: {response.status_code}")  # Debug: Check API response status
    return response.json().get("data", {}).get("value", None)

# Helper function to extract image URLs from the HTML
def extract_image_urls(description_html):
    img_tags = re.findall(r'<img [^>]*src="([^"]+)"[^>]*>', description_html)
    description_html = re.sub(r'<img [^>]*src="[^"]+"[^>]*>', '', description_html)
    return img_tags, description_html

def clean_html_and_format_text(description_html, text_frame):
    """Parses HTML and formats text correctly in PowerPoint."""
    soup = BeautifulSoup(description_html, "html.parser")

    def add_run(paragraph, text, bold=False, underline=False, font_size=10, hyperlink=None):
        """Helper to add styled text."""
        text = text.strip()
        if not text:
            return
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = "Avenir"
        if bold:
            run.font.bold = True
        if underline:
            run.font.underline = True
        if hyperlink:
            run.hyperlink.address = hyperlink

    for element in soup.contents:
        print(f"🔍 Processing element: {element}")  # Debugging

        if isinstance(element, str):
            # Handle raw text
            text = element.strip()
            if text:
                p = text_frame.add_paragraph()
                add_run(p, text)

        elif element.name in ["h1", "h2", "h3"]:
            # Add headings as bold, larger text
            p = text_frame.add_paragraph()
            add_run(p, element.get_text(strip=True), bold=True, font_size=14)

        elif element.name == "p":
            paragraph_text = element.get_text(strip=True)
            if not paragraph_text:
                continue  # Skip empty <p> tags

            p = text_frame.add_paragraph()
            
            # Check for nested formatting elements
            bold = element.find(["b", "strong"]) is not None
            underline = element.find("u") is not None
            link_tag = element.find("a")

            if link_tag and link_tag.get_text(strip=True) == paragraph_text:
                # If the entire paragraph is an <a>, treat it as a clickable link
                add_run(p, link_tag.get_text(strip=True), hyperlink=link_tag["href"])
            else:
                # Process individual parts of the paragraph
                for child in element.contents:
                    if isinstance(child, str):
                        add_run(p, child.strip(), bold=bold, underline=underline)
                    elif child.name in ["b", "strong"]:
                        add_run(p, child.get_text(strip=True), bold=True, underline=underline)
                    elif child.name == "u":
                        add_run(p, child.get_text(strip=True), underline=True, bold=bold)
                    elif child.name == "a":
                        add_run(p, child.get_text(strip=True), hyperlink=child["href"])

        elif element.name == "ul":
            # Handle unordered lists
            for li in element.find_all("li", recursive=False):
                p = text_frame.add_paragraph()
                add_run(p, f"- {li.get_text(strip=True)}")

        elif element.name == "a":
            # Standalone hyperlinks (not inside <p>)
            text = element.get_text(strip=True)
            href = element.get("href", "")
            if text and href:
                p = text_frame.add_paragraph()
                add_run(p, text, hyperlink=href)

        elif element.name == "<br/>":
            # Line breaks (only add if not redundant)
            if text_frame.paragraphs and text_frame.paragraphs[-1].text.strip():
                text_frame.add_paragraph()

    print(f"✅ Finished formatting text")


#Step 8: Insert image with maintained aspect ratio
def insert_image_with_aspect_ratio(slide, placeholder, img):
    # Get the dimensions of the placeholder
    ph_width = placeholder.width
    ph_height = placeholder.height

    # Get the dimensions of the image
    img_width, img_height = img.size

    # Calculate the aspect ratios
    ph_aspect = ph_width / ph_height
    img_aspect = img_width / img_height

    # Determine how to scale the image based on aspect ratio comparison
    if img_aspect > ph_aspect:
        # Image is wider than the placeholder (scale by width)
        new_width = ph_width
        new_height = int(ph_width / img_aspect)
    else:
        # Image is taller than the placeholder (scale by height)
        new_height = ph_height
        new_width = int(ph_height * img_aspect)

    # Calculate the position to center the image within the placeholder
    left = placeholder.left + int((ph_width - new_width) / 2)
    top = placeholder.top + int((ph_height - new_height) / 2)

    # Insert the resized image into the slide at the calculated position and size
    img_stream = BytesIO()
    img.save(img_stream, format="PNG")
    img_stream.seek(0)

    # Insert the picture into the slide
    picture = slide.shapes.add_picture(img_stream, left, top, width=new_width, height=new_height)
    
    print("Image inserted with proper aspect ratio and centered.")

# Step 9: Function to fetch Jira information for a feature
def get_jira_details(feature_id, jira_api_id):
    url = f"https://api.productboard.com/jira-integrations/{jira_api_id}/connections/{feature_id}"
    headers = {
        "accept": "application/json",
        "X-Version": "1",
        "authorization": f"Bearer {os.getenv('PRODUCTBOARD_API_TOKEN')}"
    }
    response = requests.get(url, headers=headers)
    print(f"API response for Jira connection: {response.status_code}")  # Debugging

    if response.status_code != 200:
        print(f"Failed to fetch Jira details for feature {feature_id}: {response.text}")
        return None, None

    # Extract issueKey and construct Jira URL
    jira_data = response.json()
    issue_key = jira_data.get("data", {}).get("connection", {}).get("issueKey")
    if issue_key:
        jira_url = f"https://jira.egnyte-it.com/browse/{issue_key}"
        print(f"Generated Jira URL: {jira_url}")
        return issue_key, jira_url
    else:
        print(f"No Jira issueKey found for feature {feature_id}.")
        return None, None

def list_placeholders(prs):
    for slide_layout in prs.slide_layouts:
        print(f"\nInspecting slide layout: {slide_layout.name}")
        for idx, placeholder in enumerate(slide_layout.placeholders):
            print(f"Index: {idx}, Name: {placeholder.name}, Left: {Inches(placeholder.left).inches:.2f}, "
                  f"Top: {Inches(placeholder.top).inches:.2f}, "
                  f"Width: {Inches(placeholder.width).inches:.2f}, "
                  f"Height: {Inches(placeholder.height).inches:.2f}")

def fetch_image(url):
    try:
        # Only process URLs hosted on Productboard
        if "pb-files.s3.amazonaws.com" not in url:
            print(f"Skipping non-Productboard image URL: {url}")
            return None

        img_response = requests.get(url)
        if img_response.status_code != 200:
            print(f"Failed to fetch image from {url}. Status code: {img_response.status_code}")
            return None

        # Validate the content type
        content_type = img_response.headers.get("Content-Type")
        if not content_type or not content_type.startswith("image/"):
            print(f"URL does not point to an image: {url}. Content-Type: {content_type}")
            return None

        # Open and verify the image
        img = Image.open(BytesIO(img_response.content))
        img.verify()  # Ensure it's a valid image file
        return Image.open(BytesIO(img_response.content))  # Reload the image after verification
    except Exception as e:
        print(f"Error fetching or processing image from {url}: {e}")
        return None

def add_slide(prs, feature):
    slide_layout = prs.slide_layouts[0]  # "Title and Content" layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = feature['title']
    print(f"✅ Added slide for feature: {feature['title']}")

    # Set description
    description_html = feature.get("description", "")
    images_links, cleaned_description_html = extract_image_urls(description_html)

    description_placeholder = slide.placeholders[1]  # Content placeholder
    clean_html_and_format_text(cleaned_description_html, description_placeholder.text_frame)
    print(f"✅ Set description for: {feature['title']}")

    # Insert images
    filtered_image_links = [url for url in images_links if "pb-files.s3.amazonaws.com" in url]
    if filtered_image_links:
        for idx, img_url in enumerate(filtered_image_links[:4]):  # Max 4 images per slide
            print(f"🔍 Fetching image {idx + 1} from URL: {img_url}")    
            img = fetch_image(img_url)
            if img:
                img_placeholder = slide.placeholders[10]  # Image placeholder
                insert_image_with_aspect_ratio(slide, img_placeholder, img)
                print(f"✅ Added image {idx + 1} to feature: {feature['title']}")
            else:
                print(f"⚠️ Skipped image {idx + 1} for feature: {feature['title']}")
    else:
        print(f"⚠️ No images found for feature: {feature['title']}")

    # Add requirements link
    try:
        requirements_link = feature.get("requirements_link", "Missing requirements")
        requirements_placeholder = slide.placeholders[11]  # Placeholder for requirements
        p = requirements_placeholder.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "Link to requirements"
        run.hyperlink.address = requirements_link
        print(f"✅ Set requirements link for: {feature['title']}")
    except (KeyError, IndexError):
        print(f"⚠️ No placeholder for requirements link in feature: {feature['title']}")

    # Add Jira link
    try:
        jira_name, jira_link = get_jira_details(feature["id"], JIRA_API_ID)
        if jira_name and jira_link:
            jira_placeholder = slide.placeholders[12]  # Placeholder for Jira link
            jira_paragraph = jira_placeholder.text_frame.paragraphs[0]        
            run = jira_paragraph.add_run()
            run.text = jira_name  # Display the Jira issue key
            run.hyperlink.address = jira_link
            run.font.size = Pt(10)
            run.font.name = "Avenir"
            print(f"✅ Added Jira link for: {feature['title']}")
        else:
            print(f"⚠️ No Jira link for feature: {feature['title']}")
    except (KeyError, IndexError):
        print(f"⚠️ No Jira placeholder for feature: {feature['title']}")

    # Add Productboard link
    textbox = slide.shapes.add_textbox(left=Inches(0.25), top=Inches(0.25), width=Inches(2), height=Inches(0.5))
    tb = textbox.text_frame
    pb_paragraph = tb.add_paragraph()

    run = pb_paragraph.add_run()
    run.text = "PB Link"
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.hyperlink.address = feature.get("html_link", "")
    print(f"✅ Added Productboard link for: {feature['title']}")

# Create PowerPoint slides
def create_pptx(features_data, jira_api_id, no_grouping):
    prs = Presentation("templates/corporate_template.pptx")
    list_placeholders  # Debugging tool

    if no_grouping:
        print("🚀 Running script WITHOUT initiative grouping.")
        for feature in features_data:
            add_slide(prs, feature)

    else:
        grouped_features = group_features_by_initiative(release_id)

        # Track which initiatives have been processed
        processed_initiatives = set()

        # Process each initiative and its features
        for initiative, feature_ids in grouped_features.items():
            if initiative not in processed_initiatives:
                # Add Initiative Separator Slide
                slide_layout = prs.slide_layouts[1]  # Using "Cover Page A"
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = initiative
                processed_initiatives.add(initiative)
                print(f"✅ Added initiative separator: {initiative}")

            # Add each feature under the initiative
            for feature_id in feature_ids:
                feature = next((f for f in features_data if f["id"] == feature_id), None)
                if feature:
                    add_slide(prs, feature)

    now = datetime.now()
    timestamp = now.strftime("%Y-%m-%d_%H-%M-%S")
    output_filename = f"output_presentation_{timestamp}.pptx"
    prs.save(output_filename)
    print(f"✅ Presentation saved as {output_filename}")

# Main execution
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate PowerPoint slides from Productboard release features.")
    parser.add_argument("release_id", help="The release ID for which to generate the PowerPoint slides.")
    parser.add_argument("--owner_email", help="The email of the owner to filter features by (optional).", default=None)
    parser.add_argument("--no_grouping", action="store_true", help="If set, features will NOT be grouped by initiative")

    args = parser.parse_args()
    release_id = args.release_id
    owner_email = args.owner_email
    no_grouping = args.no_grouping

    # Retrieve feature IDs
    feature_ids = get_feature_uuids(release_id)

    features_data = []
    for feature_id in feature_ids:
        feature_details = get_feature_details(feature_id)
        owner_email_from_feature = feature_details.get("data", {}).get("owner", {}).get("email", "")

        if not owner_email or owner_email_from_feature == owner_email:
            requirements_link = get_requirements_link(feature_id)

            feature_data = {
                "id": feature_id,  # Add feature ID for Jira lookup
                "title": feature_details.get("data", {}).get("name", ""),
                "description": feature_details.get("data", {}).get("description", ""),
                "images": [],
                "requirements_link": requirements_link,
                "html_link": feature_details.get("data", {}).get("links", {}).get("html", "")
            }
            features_data.append(feature_data)
        else:
            print(f"Skipping feature {feature_id} as it is not assigned to {owner_email}")

    # Pass the hardcoded Jira API ID
    create_pptx(features_data, JIRA_API_ID, no_grouping)