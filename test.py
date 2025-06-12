import argparse
import os
import re
from datetime import datetime
from io import BytesIO

import requests
import keyring
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

EXCLUDED_STATUS_IDS = [
    "402b0df5-1554-44bf-bcc4-8a11d3ca0a65",  # Candidate
    "310a6c38-3719-4d02-a6f2-9c7cdfc9a27a",  # New idea
]
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

api_token = keyring.get_password("productboard-api", "default")
if not api_token:
    raise RuntimeError("API token not found in keyring.")

HEADERS = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {api_token}"
}

TIMEFRAME_START = datetime.fromisoformat("2025-04-01")
TIMEFRAME_END = datetime.fromisoformat("2025-06-30")


def clean_html_and_format_text(description_html, text_frame):
    soup = BeautifulSoup(description_html, "html.parser")
    text_frame.clear()

    def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(11)
        run.font.name = "Avenir"
        run.font.bold = bold
        run.font.underline = underline
        if hyperlink:
            run.hyperlink.address = hyperlink

    for element in soup.contents:
        if isinstance(element, NavigableString):
            text = element.strip()
            if text:
                paragraph = text_frame.add_paragraph()
                add_run(paragraph, text)

        elif isinstance(element, Tag):
            if element.name in ['h1', 'h2', 'h3', 'h4']:
                paragraph = text_frame.add_paragraph()
                heading_text = element.get_text(strip=True)
                add_run(paragraph, heading_text, bold=True)
                paragraph.space_after = Pt(8)

            elif element.name == 'p':
                paragraph = text_frame.add_paragraph()
                for child in element.descendants:
                    if isinstance(child, NavigableString):
                        add_run(paragraph, child.strip())
                    elif child.name in ['strong', 'b']:
                        add_run(paragraph, child.get_text(strip=True), bold=True)
                    elif child.name == 'a':
                        href = child.get('href', '')
                        link_text = child.get_text(strip=True)
                        add_run(paragraph, link_text, underline=True, hyperlink=href)
                paragraph.space_after = Pt(6)

            elif element.name in ['ul', 'ol']:
                is_ordered = element.name == 'ol'
                for idx, li in enumerate(element.find_all('li', recursive=False), 1):
                    paragraph = text_frame.add_paragraph()
                    prefix = f"{idx}. " if is_ordered else "• "
                    paragraph.text = prefix
                    for li_child in li.descendants:
                        if isinstance(li_child, NavigableString):
                            paragraph.text += li_child.strip()
                        elif li_child.name in ['strong', 'b']:
                            add_run(paragraph, li_child.get_text(strip=True), bold=True)
                        elif li_child.name == 'a':
                            href = li_child.get('href', '')
                            link_text = li_child.get_text(strip=True)
                            add_run(paragraph, link_text, underline=True, hyperlink=href)
                    paragraph.space_after = Pt(4)

            elif element.name == 'br':
                text_frame.add_paragraph()


def get_all_paginated_features(url):
    all_features = []
    while url:
        response = requests.get(url, headers=HEADERS)
        data = response.json()
        all_features.extend(data.get("data", []))
        url = data.get("links", {}).get("next")
    return all_features


def get_feature_ids_by_status_id(status_id):
    url = f"https://api.productboard.com/features?status.id={status_id}"
    features = get_all_paginated_features(url)
    return set(f["id"] for f in features)


def get_feature_details(feature_id):
    url = f"https://api.productboard.com/features/{feature_id}"
    response = requests.get(url, headers=HEADERS)
    return response.json()


def get_requirements_link(feature_id):
    custom_field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={custom_field_id}&hierarchyEntity.id={feature_id}"
    response = requests.get(url, headers=HEADERS)
    return response.json().get("data", {}).get("value", None)


def create_pptx(features_data):
    prs = Presentation("templates/corporate_template.pptx")
    for feature in features_data:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = feature['title']
        desc_placeholder = slide.placeholders[1]
        clean_html_and_format_text(feature["description"], desc_placeholder.text_frame)

        pb_placeholder = slide.placeholders[13]
        pb_tf = pb_placeholder.text_frame
        pb_tf.clear()
        p_pb = pb_tf.add_paragraph()
        run_pb = p_pb.add_run()
        run_pb.text = "View in Productboard"
        run_pb.hyperlink.address = feature.get("html_link", "")
        run_pb.font.size = Pt(11)
        run_pb.font.name = "Avenir"

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    prs.save(f"output_presentation_{timestamp}.pptx")
    print(f"✅ Presentation saved as output_presentation_{timestamp}.pptx")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--owner_email", default=None, help="Filter by owner email")
    args = parser.parse_args()

    all_features = get_all_paginated_features("https://api.productboard.com/features")
    all_ids = set(f["id"] for f in all_features)

    excluded_ids = set()
    for status_id in EXCLUDED_STATUS_IDS:
        excluded_ids |= get_feature_ids_by_status_id(status_id)

    remaining_ids = list(all_ids - excluded_ids)

    features = []
    for fid in remaining_ids:
        details = get_feature_details(fid)
        data = details.get("data")
        if not isinstance(data, dict):
            continue

        tf = data.get("timeframe", {})
        start = tf.get("startDate")
        end = tf.get("endDate")
        owner = data.get("owner") or {}
        email = owner.get("email", "")

        if not (start and end):
            continue

        try:
            s_dt = datetime.fromisoformat(start)
            e_dt = datetime.fromisoformat(end)
        except Exception:
            continue

        if s_dt <= TIMEFRAME_END and e_dt >= TIMEFRAME_START:
            if not args.owner_email or args.owner_email == email:
                req = get_requirements_link(fid)
                features.append({
                    "id": fid,
                    "title": data.get("name", ""),
                    "description": data.get("description", ""),
                    "requirements_link": req,
                    "html_link": data.get("links", {}).get("html", "")
                })

    print(f"\n📊 Final features to generate slides for: {len(features)}")
    for f in features:
        print(f" - {f['title']} ({f['id']})")

    create_pptx(features)


if __name__ == "__main__":
    main()
