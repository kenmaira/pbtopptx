import argparse
import os
import re
from datetime import datetime
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
import requests
import keyring
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

TITLE_IDX = 0
DESCRIPTION_IDX = 1
IMAGE_IDX = 10
PB_LINK_IDX = 11
ID_IDX = 12
REQUIREMENTS_IDX = 13

EXCLUDED_STATUS_IDS = [
    "402b0df5-1554-44bf-bcc4-8a11d3ca0a65",  # Candidate
    "310a6c38-3719-4bff-a26c-b5d520ea1298",  # New Idea
]
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

api_token = keyring.get_password("productboard-api", "default")
if not api_token:
    raise RuntimeError("API token not found in keyring.")

HEADERS = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {api_token}"
}

TIMEFRAME_START = datetime.fromisoformat("2025-04-01")
TIMEFRAME_END = datetime.fromisoformat("2025-06-30")

#################### IMAGE-HANDLING FUNCTIONS ####################

def extract_image_urls(description_html):
    
    # Finds <img src='...'> in the HTML description and returns them, also removing those tags from the HTML so it doesn't appear as broken text.
    
    img_tags = re.findall(r'<img [^>]*src="([^"]+)"[^>]*>', description_html)
    cleaned_html = re.sub(r'<img [^>]*src="[^"]+"[^>]*>', '', description_html)
    return img_tags, cleaned_html

def fetch_image(url):
    # Fetches and verifies an image if it's hosted on pb-files.s3.amazonaws.com, if this function is not added, Emojis will be added to all features.
    try:
        if "pb-files.s3.amazonaws.com" not in url:
            print(f"Skipping non-Productboard image URL: {url}")
            return None

        resp = requests.get(url)
        if resp.status_code != 200:
            print(f"Failed to fetch image from {url}. Status code: {resp.status_code}")
            return None

        content_type = resp.headers.get("Content-Type")
        if not content_type or not content_type.startswith("image/"):
            print(f"URL is not an image: {url}")
            return None

        # Verify image integrity
        img = Image.open(BytesIO(resp.content))
        img.verify()
        return Image.open(BytesIO(resp.content))  # re-load
    except Exception as e:
        print(f"Error fetching or processing image from {url}: {e}")
        return None

def insert_image_with_aspect_ratio(slide, placeholder, img):
    
    # Inserts a PIL image into the 'placeholder' shape while preserving aspect ratio, centered within the placeholder's bounding box.

    ph_width = placeholder.width
    ph_height = placeholder.height

    img_width, img_height = img.size
    ph_aspect = ph_width / ph_height
    img_aspect = img_width / img_height

    if img_aspect > ph_aspect:
        # Width-limited
        new_width = ph_width
        new_height = int(ph_width / img_aspect)
    else:
        # Height-limited
        new_height = ph_height
        new_width = int(ph_height * img_aspect)

    left = placeholder.left + int((ph_width - new_width) / 2)
    top = placeholder.top + int((ph_height - new_height) / 2)

    img_stream = BytesIO()
    img.save(img_stream, format="PNG")
    img_stream.seek(0)

    slide.shapes.add_picture(img_stream, left, top, width=new_width, height=new_height)
    print("✅ Inserted image with aspect ratio and centered.")

#################### TEXT HELPER FUNCTIONS ####################

def safe_get(value, default="Missing requirements"):
    if value is None:
        return default
    return str(value).strip() if str(value).strip() else default

def fill_empty_text_if_needed(shape):
    """If a placeholder is truly empty, fill with a space to avoid PPTX corruption."""
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text.strip()
    if not text:
        shape.text_frame.text = " "

def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
    """Used in HTML formatting or anywhere we build runs."""
    if not text.strip():
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.font.bold = bold
    run.font.underline = underline
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

def clean_html_and_format_text(description_html, text_frame):
    # Clears the text frame, then formats the HTML into paragraphs, lists, etc.
    
    from pptx.util import Pt
    text_frame.clear()

    soup = BeautifulSoup(description_html or "", "html.parser")
    if not soup.contents:
        paragraph = text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = "The description is empty."
        run.font.size = Pt(11)
        run.font.name = "Avenir"
        return

    for element in soup.contents:
        if isinstance(element, NavigableString):
            stripped = element.strip()
            if stripped:
                p = text_frame.add_paragraph()
                add_run(p, stripped)

        elif isinstance(element, Tag):
            if element.name in ['h1', 'h2', 'h3', 'h4']:
                p = text_frame.add_paragraph()
                add_run(p, element.get_text(strip=True), bold=True)
                p.space_after = Pt(8)

            elif element.name == 'p':
                if not element.get_text(strip=True):
                    continue
                p = text_frame.add_paragraph()
                for child in element.children:
                    if isinstance(child, NavigableString):
                        add_run(p, child.strip())
                    elif child.name in ['strong', 'b']:
                        add_run(p, child.get_text(strip=True), bold=True)
                    elif child.name == 'a':
                        href = child.get('href', '')
                        link_text = child.get_text(strip=True)
                        add_run(p, link_text, underline=True, hyperlink=href)
                p.space_after = Pt(6)

            elif element.name in ['ul', 'ol']:
                is_ordered = element.name == 'ol'
                for idx, li in enumerate(element.find_all('li', recursive=False), 1):
                    p = text_frame.add_paragraph()
                    prefix = f"{idx}. " if is_ordered else "• "
                    p.text = prefix
                    for li_child in li.children:
                        if isinstance(li_child, NavigableString):
                            p.text += li_child.strip()
                        elif li_child.name in ['strong', 'b']:
                            add_run(p, li_child.get_text(strip=True), bold=True)
                        elif li_child.name == 'a':
                            href = li_child.get('href', '')
                            link_text = li_child.get_text(strip=True)
                            add_run(p, link_text, underline=True, hyperlink=href)
                    p.space_after = Pt(4)

            elif element.name == 'br':
                continue

    # Remove truly empty paragraphs
    for p in text_frame.paragraphs:
        if not p.text.strip():
            el = p._element
            el.getparent().remove(el)

    print("✅ Finished formatting text")

#################### DATA FETCH & FILTERING ####################

def get_all_paginated_features(url):
    all_features = []
    page = 1
    while url:
        print(f"📅 Fetching features page {page}: {url}")
        response = requests.get(url, headers=HEADERS)
        if response.status_code != 200:
            print(f"❌ Failed to fetch features: {response.text}")
            break
        data = response.json()
        features = data.get("data", [])
        print(f"🔹 Retrieved {len(features)} features on page {page}")
        all_features.extend(features)
        url = data.get("links", {}).get("next")
        page += 1
    print(f"📆 Total features retrieved: {len(all_features)}")
    return all_features

def get_feature_ids_by_status_id(status_id):
    url = f"https://api.productboard.com/features?status.id={status_id}"
    return set(f["id"] for f in get_all_paginated_features(url))

def get_feature_details(fid):
    url = f"https://api.productboard.com/features/{fid}"
    response = requests.get(url, headers=HEADERS)
    print(f"📄 Feature {fid} details status: {response.status_code}")
    return response.json()

def get_requirements_link(fid):
    custom_field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={custom_field_id}&hierarchyEntity.id={fid}"
    response = requests.get(url, headers=HEADERS)
    return response.json().get("data", {}).get("value", None)

def get_jira_details(fid):
    url = f"https://api.productboard.com/jira-integrations/{JIRA_API_ID}/connections/{fid}"
    response = requests.get(url, headers=HEADERS)
    if response.status_code != 200:
        print(f"⚠️ Jira link not found for {fid}: {response.status_code}")
        return None, None
    jira_data = response.json()
    issue_key = jira_data.get("data", {}).get("connection", {}).get("issueKey")
    if issue_key:
        return issue_key, f"https://jira.egnyte-it.com/browse/{issue_key}"
    return None, None

def get_placeholder_by_idx(slide, idx):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            # Fill any truly empty text with a single space
            if shape.has_text_frame and not shape.text_frame.text.strip():
                shape.text_frame.text = " "
            return shape
    return None

#################### SAFE TEXT FUNCTION ####################

def safe_clear_and_add_text(text_frame, text, hyperlink=None, size=Pt(11), font_name="Avenir"):
    """
    Clears all paragraphs in text_frame, then adds exactly one paragraph/run with 'text'.
    If text is empty, fallback to " " to avoid PPTX corruption (empty <a:p/>).
    """
    while text_frame.paragraphs:
        paragraph = text_frame.paragraphs[-1]
        for r in paragraph.runs:
            r._r.getparent().remove(r._r)  # Use r._r for the underlying <a:r>
        paragraph._element.getparent().remove(paragraph._element)

    paragraph = text_frame.add_paragraph()
    text = text.strip() if text else ""
    run = paragraph.add_run()
    run.text = text if text else " "
    run.font.size = size
    run.font.name = font_name

    # Only attach hyperlink if we have a nonempty text
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

#################### CREATE PPTX ####################

def create_pptx(features_data):
    prs = Presentation("templates/corporate_template.pptx")

    # Group features by 'initiative'
    grouped = {}
    for f in features_data:
        grouped.setdefault(f["initiative"], []).append(f)

    # For each initiative, create a cover slide and feature slides
    for initiative, items in grouped.items():
        cover_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(cover_layout)
        slide.shapes.title.text = initiative

        for feature in items:
            # 1) Extract images from the description
            all_images, cleaned_html = extract_image_urls(feature.get("description", "") or "")

            # 2) Create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = feature['title']

            # 3) Populate description
            desc_placeholder = get_placeholder_by_idx(slide, DESCRIPTION_IDX)
            if desc_placeholder and desc_placeholder.has_text_frame:
                clean_html_and_format_text(cleaned_html, desc_placeholder.text_frame)

            # 4) Insert images if we have them
            if all_images:
                img_placeholder = get_placeholder_by_idx(slide, IMAGE_IDX)
                if img_placeholder:
                    # Insert up to 4 images
                    for i, url in enumerate(all_images[:4], 1):
                        print(f"🔍 Fetching image {i} for feature: {feature['title']} => {url}")
                        pil_img = fetch_image(url)
                        if pil_img:
                            insert_image_with_aspect_ratio(slide, img_placeholder, pil_img)
                        else:
                            print(f"⚠️ Skipped invalid image at {url}")

            # 5) Productboard link
            pb_placeholder = get_placeholder_by_idx(slide, PB_LINK_IDX)
            if pb_placeholder and pb_placeholder.has_text_frame:
                link = feature.get("html_link")
                safe_clear_and_add_text(pb_placeholder.text_frame, "View in Productboard" if link else "", hyperlink=link)

            # 6) JIRA link
            jira_placeholder = get_placeholder_by_idx(slide, ID_IDX)
            if jira_placeholder and jira_placeholder.has_text_frame:
                jira_key = feature.get("jira_key")
                jira_url = feature.get("jira_url")
                if jira_key and jira_url:
                    safe_clear_and_add_text(jira_placeholder.text_frame, jira_key, hyperlink=jira_url)
                else:
                    safe_clear_and_add_text(jira_placeholder.text_frame, "No Jira Link")

            # 7) Requirements
            req_placeholder = get_placeholder_by_idx(slide, REQUIREMENTS_IDX)
            if req_placeholder and req_placeholder.has_text_frame:
                requirements_link = (feature.get("requirements_link") or "").strip()
                if requirements_link:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Link to requirements", hyperlink=requirements_link)
                else:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Missing requirements")

            # 8) Final housekeeping
            for shape in slide.placeholders:
                fill_empty_text_if_needed(shape)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    prs.save(f"output_presentation_{timestamp}.pptx")
    print(f"✅ Presentation saved as output_presentation_{timestamp}.pptx")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--owner_email", default=None, help="Filter by owner email")
    args = parser.parse_args()

    print("🔍 Retrieving all features...")
    all_features = get_all_paginated_features("https://api.productboard.com/features")
    all_ids = set(f["id"] for f in all_features)

    excluded_ids = set()
    for status_id in EXCLUDED_STATUS_IDS:
        print(f"🔎 Fetching excluded features with status {status_id}")
        excluded_ids |= get_feature_ids_by_status_id(status_id)

    remaining_ids = list(all_ids - excluded_ids)
    print(f"🧳 Remaining features after exclusion: {len(remaining_ids)}")

    initiatives = get_all_paginated_features("https://api.productboard.com/initiatives")
    initiative_map = {i["id"]: i["name"] for i in initiatives}
    feature_to_initiative = {}
    for iid in initiative_map:
        links = get_all_paginated_features(f"https://api.productboard.com/initiatives/{iid}/links/features")
        for link in links:
            fid = link.get("id")
            if fid:
                feature_to_initiative[fid] = iid

    features = []
    print("🧵 Fetching feature details in parallel...")
    with ThreadPoolExecutor(max_workers=20) as executor:
        futures = {executor.submit(get_feature_details, fid): fid for fid in remaining_ids}
        for future in as_completed(futures):
            fid = futures[future]
            try:
                details = future.result()
                data = details.get("data")
                if not isinstance(data, dict):
                    continue

                tf = data.get("timeframe", {})
                start = tf.get("startDate")
                end = tf.get("endDate")
                owner = data.get("owner") or {}
                email = owner.get("email", "")

                if not (start and end):
                    continue

                try:
                    s_dt = datetime.fromisoformat(start)
                    e_dt = datetime.fromisoformat(end)
                except Exception:
                    print(f"⚠️ Invalid date for feature {fid}: {start} - {end}")
                    continue

                if s_dt <= TIMEFRAME_END and e_dt >= TIMEFRAME_START:
                    if not args.owner_email or args.owner_email == email:
                        req = get_requirements_link(fid)
                        jira_key, jira_url = get_jira_details(fid)
                        initiative = initiative_map.get(feature_to_initiative.get(fid), "Uncategorized")
                        features.append({
                            "id": fid,
                            "title": data.get("name", ""),
                            "description": data.get("description", ""),
                            "requirements_link": req,
                            "html_link": data.get("links", {}).get("html", ""),
                            "initiative": initiative,
                            "jira_key": jira_key,
                            "jira_url": jira_url
                        })
            except Exception as e:
                print(f"❌ Error fetching feature {fid}: {e}")

    print(f"\n📊 Final features to generate slides for: {len(features)}")
    for f in features:
        print(f" - {f['title']} ({f['id']})")

    create_pptx(features)

if __name__ == "__main__":
    main()
