import argparse
import os
import re
from datetime import datetime
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed

import requests
import keyring
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

TITLE_IDX = 0
DESCRIPTION_IDX = 1
IMAGE_IDX = 10
REQUIREMENTS_IDX = 13
ID_IDX = 12
PB_LINK_IDX = 11

EXCLUDED_STATUS_IDS = [
    "402b0df5-1554-44bf-bcc4-8a11d3ca0a65",  # Candidate
    "310a6c38-3719-4d02-a6f2-9c7cdfc9a27a",  # New idea
]
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

api_token = keyring.get_password("productboard-api", "default")
if not api_token:
    raise RuntimeError("API token not found in keyring.")

HEADERS = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {api_token}"
}

TIMEFRAME_START = datetime.fromisoformat("2025-04-01")
TIMEFRAME_END = datetime.fromisoformat("2025-06-30")

def get_feature_details(fid):
    url = f"https://api.productboard.com/features/{fid}"
    response = requests.get(url, headers=HEADERS)
    print(f"📄 Feature {fid} details status: {response.status_code}")
    return response.json()

def get_requirements_link(fid):
    custom_field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={custom_field_id}&hierarchyEntity.id={fid}"
    response = requests.get(url, headers=HEADERS)
    return response.json().get("data", {}).get("value", None)

def get_jira_details(fid):
    url = f"https://api.productboard.com/jira-integrations/{JIRA_API_ID}/connections/{fid}"
    response = requests.get(url, headers=HEADERS)
    if response.status_code != 200:
        print(f"⚠️ Jira link not found for {fid}: {response.status_code}")
        return None, None
    jira_data = response.json()
    issue_key = jira_data.get("data", {}).get("connection", {}).get("issueKey")
    if issue_key:
        return issue_key, f"https://jira.egnyte-it.com/browse/{issue_key}"
    return None, None

def get_all_paginated_features(url):
    all_features = []
    page = 1
    while url:
        print(f"📅 Fetching features page {page}: {url}")
        response = requests.get(url, headers=HEADERS)
        if response.status_code != 200:
            print(f"❌ Failed to fetch features: {response.text}")
            break
        data = response.json()
        features = data.get("data", [])
        print(f"🔹 Retrieved {len(features)} features on page {page}")
        all_features.extend(features)
        url = data.get("links", {}).get("next")
        page += 1
    print(f"📆 Total features retrieved: {len(all_features)}")
    return all_features

def get_feature_ids_by_status_id(status_id):
    url = f"https://api.productboard.com/features?status.id={status_id}"
    return set(f["id"] for f in get_all_paginated_features(url))

def get_feature_details(fid):
    url = f"https://api.productboard.com/features/{fid}"
    response = requests.get(url, headers=HEADERS)
    print(f"📄 Feature {fid} details status: {response.status_code}")
    return response.json()
def safe_get(value, default="Missing requirements"):
    if value is None:
        return default
    return str(value).strip() if str(value).strip() else default

def extract_image_urls(description_html):
    img_tags = re.findall(r'<img [^>]*src="([^"]+)"[^>]*>', description_html)
    description_html = re.sub(r'<img [^>]*src="[^"]+"[^>]*>', '', description_html)
    return img_tags, description_html

def insert_image_with_aspect_ratio(slide, placeholder, img):
    ph_width = placeholder.width
    ph_height = placeholder.height
    img_width, img_height = img.size
    ph_aspect = ph_width / ph_height
    img_aspect = img_width / img_height
    if img_aspect > ph_aspect:
        new_width = ph_width
        new_height = int(ph_width / img_aspect)
    else:
        new_height = ph_height
        new_width = int(ph_height * img_aspect)
    left = placeholder.left + int((ph_width - new_width) / 2)
    top = placeholder.top + int((ph_height - new_height) / 2)
    img_stream = BytesIO()
    img.save(img_stream, format="PNG")
    img_stream.seek(0)
    slide.shapes.add_picture(img_stream, left, top, width=new_width, height=new_height)

def fetch_image(url):
    try:
        if "pb-files.s3.amazonaws.com" not in url:
            print(f"Skipping non-Productboard image URL: {url}")
            return None
        img_response = requests.get(url)
        if img_response.status_code != 200:
            print(f"Failed to fetch image from {url}. Status code: {img_response.status_code}")
            return None
        content_type = img_response.headers.get("Content-Type")
        if not content_type or not content_type.startswith("image/"):
            print(f"URL does not point to an image: {url}. Content-Type: {content_type}")
            return None
        img = Image.open(BytesIO(img_response.content))
        img.verify()
        return Image.open(BytesIO(img_response.content))
    except Exception as e:
        print(f"Error fetching or processing image from {url}: {e}")
        return None

def fill_empty_text_if_needed(shape):
    if not shape.has_text_frame:
        return
    if not shape.text_frame.text.strip():
        shape.text_frame.text = " "

def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
    if not text.strip():
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.font.bold = bold
    run.font.underline = underline
    if hyperlink:
        run.hyperlink.address = hyperlink

def render_description(description_html, text_frame):
    soup = BeautifulSoup(description_html or "", 'html.parser')
    if not soup.contents:
        paragraph = text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = "The description is empty."
        run.font.size = Pt(11)
        run.font.name = "Avenir"
        return
    for element in soup.contents:
        if isinstance(element, NavigableString):
            if element.strip():
                p = text_frame.add_paragraph()
                add_run(p, element.strip())
        elif isinstance(element, Tag):
            if element.name == 'p':
                p = text_frame.add_paragraph()
                for child in element.children:
                    if isinstance(child, NavigableString):
                        add_run(p, child.strip())
                    elif child.name in ['strong', 'b']:
                        add_run(p, child.get_text(strip=True), bold=True)
                    elif child.name == 'a':
                        href = child.get('href', '')
                        link_text = child.get_text(strip=True)
                        add_run(p, link_text, underline=True, hyperlink=href)
                p.space_after = Pt(6)
            elif element.name in ['ul', 'ol']:
                is_ordered = element.name == 'ol'
                for idx, li in enumerate(element.find_all('li', recursive=False), 1):
                    p = text_frame.add_paragraph()
                    prefix = f"{idx}. " if is_ordered else "• "
                    p.text = prefix
                    for li_child in li.children:
                        if isinstance(li_child, NavigableString):
                            p.text += li_child.strip()
                        elif li_child.name in ['strong', 'b']:
                            add_run(p, li_child.get_text(strip=True), bold=True)
                        elif li_child.name == 'a':
                            href = li_child.get('href', '')
                            link_text = li_child.get_text(strip=True)
                            add_run(p, link_text, underline=True, hyperlink=href)
                    p.space_after = Pt(4)
    for p in text_frame.paragraphs:
        if not p.text.strip():
            el = p._element
            el.getparent().remove(el)

def create_pptx(features):
    prs = Presentation("templates/corporate_template.pptx")
    features.sort(key=lambda x: x.get("initiative", "Uncategorized"))
    current_initiative = None
    for feature in features:
        initiative = feature.get("initiative", "Uncategorized")
        if initiative != current_initiative:
            sep_slide_layout = prs.slide_layouts[1]
            sep_slide = prs.slides.add_slide(sep_slide_layout)
            title_shape = sep_slide.shapes.title or sep_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
            title_shape.text = initiative
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Avenir"
                    run.font.size = Pt(28)
            current_initiative = initiative

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.placeholders[TITLE_IDX].text = feature.get("title", "Missing title")
        desc_html = feature.get("description", "")
        image_urls, cleaned_desc = extract_image_urls(desc_html)
        desc_shape = slide.placeholders[DESCRIPTION_IDX]
        desc_shape.text = ""
        render_description(cleaned_desc, desc_shape.text_frame)

        if IMAGE_IDX < len(slide.placeholders):
            image_shape = slide.placeholders[IMAGE_IDX]
            image_shape.text = ""
            for url in image_urls:
                img = fetch_image(url)
                if img:
                    insert_image_with_aspect_ratio(slide, image_shape, img)
                    break

        if REQUIREMENTS_IDX < len(slide.placeholders):
            requirements_shape = slide.placeholders[REQUIREMENTS_IDX]
            requirements_shape.text_frame.clear()
            p = requirements_shape.text_frame.add_paragraph()
            add_run(p, "Link to requirements", underline=True, hyperlink=safe_get(feature.get("requirements_link")))

        if ID_IDX < len(slide.placeholders):
            id_shape = slide.placeholders[ID_IDX]
            id_shape.text_frame.clear()
            p = id_shape.text_frame.add_paragraph()
            ticket_id = feature.get("jira_key", "")
            if ticket_id:
                add_run(p, ticket_id, underline=True, hyperlink=feature.get("jira_url"))
            else:
                fill_empty_text_if_needed(id_shape)

        if PB_LINK_IDX < len(slide.placeholders):
            pb_link_shape = slide.placeholders[PB_LINK_IDX]
            pb_link_shape.text_frame.clear()
            p = pb_link_shape.text_frame.add_paragraph()
            pb_url = feature.get("html_link", "")
            if pb_url:
                add_run(p, "View in Productboard", underline=True, hyperlink=pb_url)
            else:
                add_run(p, "Missing PB link")

        for shape in slide.placeholders:
            fill_empty_text_if_needed(shape)

    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    prs.save(f"output_presentation_{timestamp}.pptx")
    print(f"✅ Presentation saved as output_presentation_{timestamp}.pptx")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--owner_email", default=None, help="Filter by owner email")
    args = parser.parse_args()

    print("🔍 Retrieving all features...")
    all_features = get_all_paginated_features("https://api.productboard.com/features")
    all_ids = set(f["id"] for f in all_features)

    excluded_ids = set()
    for status_id in EXCLUDED_STATUS_IDS:
        print(f"🔎 Fetching excluded features with status {status_id}")
        excluded_ids |= get_feature_ids_by_status_id(status_id)

    remaining_ids = list(all_ids - excluded_ids)
    print(f"🧳 Remaining features after exclusion: {len(remaining_ids)}")

    initiatives = get_all_paginated_features("https://api.productboard.com/initiatives")
    initiative_map = {i["id"]: i["name"] for i in initiatives}
    feature_to_initiative = {}
    for iid in initiative_map:
        links = get_all_paginated_features(f"https://api.productboard.com/initiatives/{iid}/links/features")
        for link in links:
            fid = link.get("id")
            if fid:
                feature_to_initiative[fid] = iid

    features = []
    print("🧵 Fetching feature details in parallel...")
    with ThreadPoolExecutor(max_workers=20) as executor:
        futures = {executor.submit(get_feature_details, fid): fid for fid in remaining_ids}
        for future in as_completed(futures):
            fid = futures[future]
            try:
                details = future.result()
                data = details.get("data")
                if not isinstance(data, dict):
                    continue

                tf = data.get("timeframe", {})
                start = tf.get("startDate")
                end = tf.get("endDate")
                owner = data.get("owner") or {}
                email = owner.get("email", "")

                if not (start and end):
                    continue

                try:
                    s_dt = datetime.fromisoformat(start)
                    e_dt = datetime.fromisoformat(end)
                except Exception:
                    print(f"⚠️ Invalid date for feature {fid}: {start} - {end}")
                    continue

                if s_dt <= TIMEFRAME_END and e_dt >= TIMEFRAME_START:
                    if not args.owner_email or args.owner_email == email:
                        req = get_requirements_link(fid)
                        jira_key, jira_url = get_jira_details(fid)
                        initiative = initiative_map.get(feature_to_initiative.get(fid), "Uncategorized")
                        features.append({
                            "id": fid,
                            "title": data.get("name", ""),
                            "description": data.get("description", ""),
                            "requirements_link": req,
                            "html_link": data.get("links", {}).get("html", ""),
                            "initiative": initiative,
                            "jira_key": jira_key,
                            "jira_url": jira_url
                        })
            except Exception as e:
                print(f"❌ Error fetching feature {fid}: {e}")

    print(f"\n📊 Final features to generate slides for: {len(features)}")
    for f in features:
        print(f" - {f['title']} ({f['id']})")

    create_pptx(features)

if __name__ == "__main__":
    main()
