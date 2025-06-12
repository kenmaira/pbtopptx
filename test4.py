import argparse
import os
import re
from datetime import datetime
from io import BytesIO
from concurrent.futures import ThreadPoolExecutor, as_completed
import requests
import keyring
from bs4 import BeautifulSoup, NavigableString, Tag
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

TITLE_IDX = 0
DESCRIPTION_IDX = 1
IMAGE_IDX = 10
PB_LINK_IDX = 11
ID_IDX = 12
REQUIREMENTS_IDX = 13

EXCLUDED_STATUS_IDS = [
    "402b0df5-1554-44bf-bcc4-8a11d3ca0a65",  # Candidate
    "310a6c38-3719-4d02-a6f2-9c7cdfc9a27a",  # New idea
]
JIRA_API_ID = "155d80cb-8b4c-4bff-a26c-b5d520ea1298"

api_token = keyring.get_password("productboard-api", "default")
if not api_token:
    raise RuntimeError("API token not found in keyring.")

HEADERS = {
    "accept": "application/json",
    "X-Version": "1",
    "authorization": f"Bearer {api_token}"
}

TIMEFRAME_START = datetime.fromisoformat("2025-04-01")
TIMEFRAME_END = datetime.fromisoformat("2025-06-30")

def safe_get(value, default="Missing requirements"):
    if value is None:
        return default
    return str(value).strip() if str(value).strip() else default

def fill_empty_text_if_needed(shape):
    """If a placeholder is truly empty, fill with a space to avoid PPTX corruption."""
    if not shape.has_text_frame:
        return
    text = shape.text_frame.text.strip()
    if not text:
        shape.text_frame.text = " "

def add_run(paragraph, text, bold=False, underline=False, hyperlink=None):
    """Used in HTML description formatting and anywhere else we build runs manually."""
    if not text.strip():
        return
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.name = "Avenir"
    run.font.bold = bold
    run.font.underline = underline
    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

def clean_html_and_format_text(description_html, text_frame):
    """Clears the text frame, then formats HTML description with paragraphs/lists/links."""
    soup = BeautifulSoup(description_html or "", "html.parser")
    text_frame.clear()  # This can safely remove content (pptx expects a paragraph after this)

    if not soup.contents:
        paragraph = text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = "The description is empty."
        run.font.size = Pt(11)
        run.font.name = "Avenir"
        return

    for element in soup.contents:
        if isinstance(element, NavigableString):
            stripped = element.strip()
            if stripped:
                p = text_frame.add_paragraph()
                add_run(p, stripped)

        elif isinstance(element, Tag):
            if element.name in ['h1', 'h2', 'h3', 'h4']:
                p = text_frame.add_paragraph()
                add_run(p, element.get_text(strip=True), bold=True)
                p.space_after = Pt(8)

            elif element.name == 'p':
                if not element.get_text(strip=True):
                    continue
                p = text_frame.add_paragraph()
                for child in element.children:
                    if isinstance(child, NavigableString):
                        add_run(p, child.strip())
                    elif child.name in ['strong', 'b']:
                        add_run(p, child.get_text(strip=True), bold=True)
                    elif child.name == 'a':
                        href = child.get('href', '')
                        link_text = child.get_text(strip=True)
                        add_run(p, link_text, underline=True, hyperlink=href)
                p.space_after = Pt(6)

            elif element.name in ['ul', 'ol']:
                is_ordered = element.name == 'ol'
                for idx, li in enumerate(element.find_all('li', recursive=False), 1):
                    p = text_frame.add_paragraph()
                    prefix = f"{idx}. " if is_ordered else "• "
                    p.text = prefix
                    for li_child in li.children:
                        if isinstance(li_child, NavigableString):
                            p.text += li_child.strip()
                        elif li_child.name in ['strong', 'b']:
                            add_run(p, li_child.get_text(strip=True), bold=True)
                        elif li_child.name == 'a':
                            href = li_child.get('href', '')
                            link_text = li_child.get_text(strip=True)
                            add_run(p, link_text, underline=True, hyperlink=href)
                    p.space_after = Pt(4)

            elif element.name == 'br':
                continue

    # Remove any truly empty paragraphs
    for p in text_frame.paragraphs:
        if not p.text.strip():
            el = p._element
            el.getparent().remove(el)

    print("✅ Finished formatting text")

def get_all_paginated_features(url):
    all_features = []
    page = 1
    while url:
        print(f"📅 Fetching features page {page}: {url}")
        response = requests.get(url, headers=HEADERS)
        if response.status_code != 200:
            print(f"❌ Failed to fetch features: {response.text}")
            break
        data = response.json()
        features = data.get("data", [])
        print(f"🔹 Retrieved {len(features)} features on page {page}")
        all_features.extend(features)
        url = data.get("links", {}).get("next")
        page += 1
    print(f"📆 Total features retrieved: {len(all_features)}")
    return all_features

def get_feature_ids_by_status_id(status_id):
    url = f"https://api.productboard.com/features?status.id={status_id}"
    return set(f["id"] for f in get_all_paginated_features(url))

def get_feature_details(fid):
    url = f"https://api.productboard.com/features/{fid}"
    response = requests.get(url, headers=HEADERS)
    print(f"📄 Feature {fid} details status: {response.status_code}")
    return response.json()

def get_requirements_link(fid):
    custom_field_id = "52ae58e7-6417-4898-956b-bd74d4e87502"
    url = f"https://api.productboard.com/hierarchy-entities/custom-fields-values/value?customField.id={custom_field_id}&hierarchyEntity.id={fid}"
    response = requests.get(url, headers=HEADERS)
    return response.json().get("data", {}).get("value", None)

def get_jira_details(fid):
    url = f"https://api.productboard.com/jira-integrations/{JIRA_API_ID}/connections/{fid}"
    response = requests.get(url, headers=HEADERS)
    if response.status_code != 200:
        print(f"⚠️ Jira link not found for {fid}: {response.status_code}")
        return None, None
    jira_data = response.json()
    issue_key = jira_data.get("data", {}).get("connection", {}).get("issueKey")
    if issue_key:
        return issue_key, f"https://jira.egnyte-it.com/browse/{issue_key}"
    return None, None

def get_placeholder_by_idx(slide, idx):
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            if shape.has_text_frame and not shape.text_frame.text.strip():
                shape.text_frame.text = " "
            return shape
    return None

def safe_clear_and_add_text(text_frame, text, hyperlink=None, size=Pt(11), font_name="Avenir"):
    # Remove all existing paragraphs/runs
    while text_frame.paragraphs:
        paragraph = text_frame.paragraphs[-1]
        for r in paragraph.runs:
            # Use r._r instead of r._element
            r._r.getparent().remove(r._r)
        paragraph._element.getparent().remove(paragraph._element)

    paragraph = text_frame.add_paragraph()
    text = text.strip() if text else ""
    run = paragraph.add_run()
    run.text = text if text else " "
    run.font.size = size
    run.font.name = font_name

    if hyperlink and text.strip():
        run.hyperlink.address = hyperlink

def create_pptx(features_data):
    prs = Presentation("templates/corporate_template.pptx")

    # Group features by 'initiative' so we can create a "cover slide" for each
    grouped = {}
    for f in features_data:
        grouped.setdefault(f["initiative"], []).append(f)

    # Create slides
    for initiative, items in grouped.items():
        # Cover slide for this initiative (use layout index=1)
        cover_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(cover_layout)
        slide.shapes.title.text = initiative

        # Now create slides for each feature in this initiative
        for feature in items:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = feature['title']

            # Description placeholder (idx=1)
            desc_placeholder = get_placeholder_by_idx(slide, 1)
            if desc_placeholder and desc_placeholder.has_text_frame:
                clean_html_and_format_text(feature.get("description", ""), desc_placeholder.text_frame)

            # Productboard link placeholder (idx=13 or PB_LINK_IDX=11, verify your indexing!)
            pb_placeholder = get_placeholder_by_idx(slide, PB_LINK_IDX)
            if pb_placeholder and pb_placeholder.has_text_frame:
                link = feature.get("html_link")
                # Insert "View in Productboard" or blank
                safe_clear_and_add_text(pb_placeholder.text_frame, "View in Productboard" if link else "", hyperlink=link)
            else:
                print(f"⚠️ Placeholder {PB_LINK_IDX} missing or invalid for PB link on slide: {feature['title']}")

            # JIRA link placeholder (idx=12 or ID_IDX=12)
            jira_placeholder = get_placeholder_by_idx(slide, ID_IDX)
            if jira_placeholder and jira_placeholder.has_text_frame:
                jira_key = feature.get("jira_key")
                jira_url = feature.get("jira_url")
                if jira_key and jira_url:
                    safe_clear_and_add_text(jira_placeholder.text_frame, jira_key, hyperlink=jira_url)
                else:
                    safe_clear_and_add_text(jira_placeholder.text_frame, "No Jira Link")
            else:
                print(f"⚠️ Placeholder {ID_IDX} missing or invalid for JIRA on slide: {feature['title']}")

            # Requirements placeholder (idx=11 or REQUIREMENTS_IDX=13)
            req_placeholder = get_placeholder_by_idx(slide, REQUIREMENTS_IDX)
            if req_placeholder and req_placeholder.has_text_frame:
                requirements_link = feature.get("requirements_link") or ""
                if requirements_link.strip():
                    safe_clear_and_add_text(req_placeholder.text_frame, "Link to requirements", hyperlink=requirements_link.strip())
                else:
                    safe_clear_and_add_text(req_placeholder.text_frame, "Missing requirements")
            else:
                print(f"⚠️ Placeholder {REQUIREMENTS_IDX} missing or invalid on slide: {feature['title']}")

            # Final pass to fill truly empty placeholders with a space
            for shape in slide.placeholders:
                fill_empty_text_if_needed(shape)

    # Save final
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    prs.save(f"output_presentation_{timestamp}.pptx")
    print(f"✅ Presentation saved as output_presentation_{timestamp}.pptx")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--owner_email", default=None, help="Filter by owner email")
    args = parser.parse_args()

    print("🔍 Retrieving all features...")
    all_features = get_all_paginated_features("https://api.productboard.com/features")
    all_ids = set(f["id"] for f in all_features)

    excluded_ids = set()
    for status_id in EXCLUDED_STATUS_IDS:
        print(f"🔎 Fetching excluded features with status {status_id}")
        excluded_ids |= get_feature_ids_by_status_id(status_id)

    remaining_ids = list(all_ids - excluded_ids)
    print(f"🧳 Remaining features after exclusion: {len(remaining_ids)}")

    initiatives = get_all_paginated_features("https://api.productboard.com/initiatives")
    initiative_map = {i["id"]: i["name"] for i in initiatives}
    feature_to_initiative = {}
    for iid in initiative_map:
        links = get_all_paginated_features(f"https://api.productboard.com/initiatives/{iid}/links/features")
        for link in links:
            fid = link.get("id")
            if fid:
                feature_to_initiative[fid] = iid

    features = []
    print("🧵 Fetching feature details in parallel...")
    with ThreadPoolExecutor(max_workers=20) as executor:
        futures = {executor.submit(get_feature_details, fid): fid for fid in remaining_ids}
        for future in as_completed(futures):
            fid = futures[future]
            try:
                details = future.result()
                data = details.get("data")
                if not isinstance(data, dict):
                    continue

                tf = data.get("timeframe", {})
                start = tf.get("startDate")
                end = tf.get("endDate")
                owner = data.get("owner") or {}
                email = owner.get("email", "")

                if not (start and end):
                    continue

                try:
                    s_dt = datetime.fromisoformat(start)
                    e_dt = datetime.fromisoformat(end)
                except Exception:
                    print(f"⚠️ Invalid date for feature {fid}: {start} - {end}")
                    continue

                if s_dt <= TIMEFRAME_END and e_dt >= TIMEFRAME_START:
                    if not args.owner_email or args.owner_email == email:
                        req = get_requirements_link(fid)
                        jira_key, jira_url = get_jira_details(fid)
                        initiative = initiative_map.get(feature_to_initiative.get(fid), "Uncategorized")
                        features.append({
                            "id": fid,
                            "title": data.get("name", ""),
                            "description": data.get("description", ""),
                            "requirements_link": req,
                            "html_link": data.get("links", {}).get("html", ""),
                            "initiative": initiative,
                            "jira_key": jira_key,
                            "jira_url": jira_url
                        })
            except Exception as e:
                print(f"❌ Error fetching feature {fid}: {e}")

    print(f"\n📊 Final features to generate slides for: {len(features)}")
    for f in features:
        print(f" - {f['title']} ({f['id']})")

    create_pptx(features)

if __name__ == "__main__":
    main()
